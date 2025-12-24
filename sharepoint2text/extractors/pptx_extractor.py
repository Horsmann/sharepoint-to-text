import io
import logging
from datetime import datetime

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger(__name__)


def _dt_to_iso(dt: datetime | None) -> str | None:
    return dt.isoformat() if dt else None


def read_pptx(file_like: io.BytesIO) -> dict:
    logger.debug("Reading pptx")
    prs = Presentation(file_like)

    cp = prs.core_properties
    metadata = {
        "title": cp.title,
        "subject": cp.subject,
        "author": cp.author,
        "last_modified_by": cp.last_modified_by,
        "created": _dt_to_iso(cp.created),
        "modified": _dt_to_iso(cp.modified),
        "keywords": cp.keywords,
        "comments": cp.comments,
        "category": cp.category,
        "revision": cp.revision,
    }

    slides_result: list[dict] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        logger.debug(f"Processing slide [{slide_index}]")

        slide_data = {
            "slide_number": slide_index,
            "title": None,
            "footer": None,
            "content_placeholders": [],
            "other_textboxes": [],
            "images": [],
        }

        image_counter = 0

        for shape in slide.shapes:
            # ---------------------------
            # Image extraction
            # ---------------------------
            if shape.shape_type == shape.shape_type.PICTURE:
                try:
                    image = shape.image
                    image_counter += 1

                    slide_data["images"].append(
                        {
                            "image_index": image_counter,
                            "filename": image.filename,
                            "content_type": image.content_type,
                            "size_bytes": len(image.blob),
                            "blob": image.blob,  # raw binary bytes
                        }
                    )
                except Exception as e:
                    logger.error(e)
                    logger.exception(f"Failed to extract image on slide {slide_index}")
                continue

            # ---------------------------
            # Text extraction
            # ---------------------------
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            if shape.is_placeholder:
                ptype = shape.placeholder_format.type

                if ptype in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.VERTICAL_TITLE,
                ):
                    slide_data["title"] = text

                elif ptype == PP_PLACEHOLDER.FOOTER:
                    slide_data["footer"] = text

                elif ptype in (
                    PP_PLACEHOLDER.BODY,
                    PP_PLACEHOLDER.SUBTITLE,
                    PP_PLACEHOLDER.OBJECT,
                    PP_PLACEHOLDER.VERTICAL_BODY,
                    PP_PLACEHOLDER.VERTICAL_OBJECT,
                    PP_PLACEHOLDER.TABLE,
                ):
                    slide_data["content_placeholders"].append(text)

                else:
                    slide_data["other_textboxes"].append(text)
            else:
                slide_data["other_textboxes"].append(text)

        slides_result.append(slide_data)

    return {
        "metadata": metadata,
        "slides": slides_result,
    }
